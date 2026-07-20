import io
import os

import openpyxl
import pdfplumber
from docx import Document
from dotenv import load_dotenv
from langchain_core.runnables import RunnableConfig
from langchain_core.tools import tool
from pptx import Presentation

from tools.sandbox_utils import get_sandbox_sync

load_dotenv()

sandbox_temp_dir = "/tmp"


def _find_cjk_font() -> str | None:
    """查找系统中可用的 CJK/Unicode 字体。"""
    import platform

    candidates = []
    if platform.system() == "Windows":
        candidates = [
            "C:/Windows/Fonts/msyh.ttc",       # 微软雅黑
            "C:/Windows/Fonts/simsun.ttc",      # 宋体
            "C:/Windows/Fonts/arial.ttf",
        ]
    else:
        candidates = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
            "/usr/share/fonts/truetype/arphic/uming.ttc",
        ]

    for path in candidates:
        if os.path.isfile(path):
            return path
    return None


def _parse_file_bytes(content: bytes, ext: str) -> str:
    """根据文件扩展名解析文件内容为纯文本。

    Args:
        content: 文件的原始字节内容。
        ext: 文件扩展名（含点号，如 ".pdf"）。

    Returns:
        解析后的纯文本。
    """
    if ext in (".md", ".txt", ".csv", ".json", ".xml", ".py", ".js", ".ts", ".html", ".css", ".yaml", ".yml"):
        return content.decode("utf-8")

    if ext == ".pdf":
        parts: list[str] = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text)
        return "\n".join(parts)

    if ext in (".xlsx", ".xls"):
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)

    if ext == ".docx":
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    if ext == ".pptx":
        prs = Presentation(io.BytesIO(content))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)

    # 未知格式尝试按 UTF-8 文本处理
    return content.decode("utf-8")


def _generate_pdf_bytes(text: str) -> bytes:
    """使用 fpdf2 生成 PDF 文件，支持自动换行、分页和 CJK 字符。"""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    cjk_font = _find_cjk_font()
    if cjk_font:
        pdf.add_font("uni", "", cjk_font)
        pdf.add_page()
        pdf.set_font("uni", size=11)
    else:
        pdf.add_page()
        pdf.set_font("Courier", size=10)

    for line in text.split("\n"):
        if line == "":
            pdf.ln(5)
        else:
            pdf.multi_cell(w=0, h=5, text=line)

    return pdf.output()


@tool
def read_file(file_path: str, config: RunnableConfig) -> str:
    """读取 sandbox 中的文件内容，支持 .pdf / .md / .docx / .xlsx / .pptx / .txt 等格式。

    Args:
        file_path: 文件路径（sandbox 中的绝对路径或相对路径）。
        config: LangChain 运行时配置（自动注入），用于获取 sandbox。

    Returns:
        文件内容（纯文本）。
    """
    try:
        sandbox = get_sandbox_sync(config)
        content = sandbox.files.read_bytes(file_path)

        _, ext = os.path.splitext(file_path)
        return _parse_file_bytes(content, ext.lower())
    except Exception as e:
        return f"读取文件失败: {e}"


@tool
def save_to_markdown(text: str, file_name: str, config: RunnableConfig) -> str:
    """将 text 保存为 markdown 文件到 sandbox 的 /tmp 目录。

    Args:
        text: 需要保存的文本。
        file_name: 保存的文件名（不含路径，如 "report.md"）。

    Returns:
        保存成功返回文件路径，失败返回错误信息。
    """
    try:
        sandbox = get_sandbox_sync(config)
        if not file_name.endswith(".md"):
            file_name += ".md"
        remote_path = f"{sandbox_temp_dir}/{file_name}"
        sandbox.files.write_file(remote_path, text.encode("utf-8"))
        return f"文件已保存至 {remote_path}"
    except Exception as e:
        return f"保存 markdown 文件失败: {e}"


@tool
def save_to_pdf(text: str, file_name: str, config: RunnableConfig) -> str:
    """将 text 保存为 PDF 文件到 sandbox 的 /tmp 目录。

    Args:
        text: 需要保存的文本内容。
        file_name: 保存的文件名（不含路径，如 "report.pdf"）。

    Returns:
        保存成功返回文件路径，失败返回错误信息。
    """
    try:
        sandbox = get_sandbox_sync(config)
        if not file_name.endswith(".pdf"):
            file_name += ".pdf"
        remote_path = f"{sandbox_temp_dir}/{file_name}"
        pdf_bytes = _generate_pdf_bytes(text)
        sandbox.files.write_file(remote_path, pdf_bytes)
        return f"文件已保存至 {remote_path}"
    except Exception as e:
        return f"保存 PDF 文件失败: {e}"


import os as _os

_API_BASE = _os.getenv("MILO_API_BASE", "http://localhost:8000")


@tool
def generate_download_url_from_sandbox(file_path: str, config: RunnableConfig) -> str:
    """Generate a download URL for a file in the sandbox.

    Returns a URL pointing to the MiloAgent file download API,
    which reads the file directly from the sandbox via OpenSandbox API.

    Verifies the file exists before returning the URL.

    Args:
        file_path: Absolute path to the file in the sandbox, e.g. "/tmp/report.pdf".
        config: LangChain runtime config (auto-injected).
    Returns:
        Download URL for the file, or error message if file not found.
    """
    import urllib.parse

    try:
        sandbox = get_sandbox_sync(config)
        # Verify file exists by listing or attempting a stat
        check = sandbox.commands.run(f"test -f {file_path} && echo 'exists' || echo 'missing'")
        stdout = check.logs.stdout[0].text if check.logs.stdout else ""
        if "missing" in stdout:
            # Try to list /tmp contents to help debugging
            ls = sandbox.commands.run(f"ls -la {file_path} 2>&1 || ls -la /tmp/ 2>&1 | head -20")
            ls_out = ls.logs.stdout[0].text if ls.logs.stdout else ""
            return f"文件不存在: {file_path}\n\n/tmp/ 目录内容:\n{ls_out}"
    except Exception as e:
        return f"无法验证文件: {e}"

    cfg = config.get("configurable", {})
    thread_id = cfg.get("thread_id", "")
    token = cfg.get("token", "")
    encoded = urllib.parse.quote(file_path, safe="")
    url = f"{_API_BASE}/api/files/download?thread_id={thread_id}&file_path={encoded}"
    if token:
        url += f"&token={token}"
    return url


# 获取文件相关的tools
def get_file_tools():
    return [read_file, save_to_pdf, save_to_markdown, generate_download_url_from_sandbox]
